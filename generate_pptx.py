import os
import sys

def check_and_install_dependencies():
    try:
        import pptx
    except ImportError:
        print("python-pptx not found. Installing...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx

check_and_install_dependencies()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    BG_COLOR = RGBColor(10, 15, 30)        # Deep Navy/Black
    CARD_BG = RGBColor(20, 27, 50)        # Dark Slate Card
    TEXT_MAIN = RGBColor(245, 247, 250)    # Crisp White/Off-white
    TEXT_MUTED = RGBColor(150, 165, 185)  # Muted Blue-Grey
    ACCENT_ORANGE = RGBColor(249, 115, 22) # SITCOE Orange Accent
    ACCENT_BLUE = RGBColor(14, 165, 233)   # Tech Sky Blue
    
    # Fonts
    FONT_TITLE = "Arial" # Fallback to system font that looks clean and modern
    FONT_BODY = "Arial"
    
    def apply_solid_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category="SMART CAMPUS NAVIGATION"):
        # Header category tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = FONT_TITLE
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_ORANGE
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        
        # Accent decorative line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.color.rgb = ACCENT_BLUE

    def add_footer(slide, current_slide, total_slides=9):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"SITCOE Smart Campus Navigation System  |  Slide {current_slide} of {total_slides}"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide_layout = prs.slide_layouts[6] # Blank Layout
    slide1 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide1, BG_COLOR)
    
    # Visual grid elements (simulating modern UI background)
    line1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.06), Inches(3.2))
    line1.fill.solid()
    line1.fill.fore_color.rgb = ACCENT_ORANGE
    line1.line.fill.background()
    
    # Title Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_trust = tf1.paragraphs[0]
    p_trust.text = "Shri Shamrao Patil Yadravkar Educational & Charitable Trust's"
    p_trust.font.name = FONT_BODY
    p_trust.font.size = Pt(13)
    p_trust.font.color.rgb = ACCENT_BLUE
    p_trust.space_after = Pt(8)
    
    p_coll = tf1.add_paragraph()
    p_coll.text = "SHARAD INSTITUTE OF TECHNOLOGY\nCOLLEGE OF ENGINEERING"
    p_coll.font.name = FONT_TITLE
    p_coll.font.size = Pt(36)
    p_coll.font.bold = True
    p_coll.font.color.rgb = TEXT_MAIN
    p_coll.space_after = Pt(16)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Smart Campus Navigation System"
    p_sub.font.name = FONT_TITLE
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ACCENT_ORANGE
    
    # Info Box / Presenter Detail
    pres_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(11.0), Inches(1.2))
    tf_pres = pres_box.text_frame
    tf_pres.word_wrap = True
    
    p_pres = tf_pres.paragraphs[0]
    p_pres.text = "Interactive presentation on indoor routing, department directory & intelligent search"
    p_pres.font.name = FONT_BODY
    p_pres.font.size = Pt(11)
    p_pres.font.color.rgb = TEXT_MUTED
    
    p_by = tf_pres.add_paragraph()
    p_by.text = "Designed for SITCOE Students, Faculty & Visitors"
    p_by.font.name = FONT_BODY
    p_by.font.size = Pt(11)
    p_by.font.bold = True
    p_by.font.color.rgb = TEXT_MAIN
    p_by.space_before = Pt(6)

    add_footer(slide1, 1)

    # ==================== SLIDE 2: THE PROBLEM ====================
    slide2 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide2, BG_COLOR)
    add_header(slide2, "The Problem")
    
    # Subtitle or hook
    hook_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.5))
    tf_hook = hook_box.text_frame
    p_hook = tf_hook.paragraphs[0]
    p_hook.text = "Campus navigation can be stressful and inefficient for both newcomers and regulars."
    p_hook.font.name = FONT_BODY
    p_hook.font.size = Pt(14)
    p_hook.font.bold = True
    p_hook.font.color.rgb = ACCENT_BLUE

    # Card 1: Time Loss
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(3.6), Inches(3.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BG
    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_right = tf_c1.margin_top = tf_c1.margin_bottom = Inches(0.3)
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "⏰ Time Loss & Delay"
    p_c1_t.font.name = FONT_TITLE
    p_c1_t.font.size = Pt(18)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_ORANGE
    p_c1_t.space_after = Pt(14)
    p_c1_d = tf_c1.add_paragraph()
    p_c1_d.text = "Students, external examiners, and guests lose precious minutes trying to locate specific exam centers, labs, seminars, or classrooms across multiple blocks."
    p_c1_d.font.name = FONT_BODY
    p_c1_d.font.size = Pt(12)
    p_c1_d.font.color.rgb = TEXT_MUTED

    # Card 2: Campus Expansion
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(2.6), Inches(3.6), Inches(3.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = CARD_BG
    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_right = tf_c2.margin_top = tf_c2.margin_bottom = Inches(0.3)
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🏫 Large Campus Layouts"
    p_c2_t.font.name = FONT_TITLE
    p_c2_t.font.size = Pt(18)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_ORANGE
    p_c2_t.space_after = Pt(14)
    p_c2_d = tf_c2.add_paragraph()
    p_c2_d.text = "SITCOE features multiple wings, departments, faculty offices, and laboratory zones. The physical layout makes manual directions confusing to explain and hard to recall."
    p_c2_d.font.name = FONT_BODY
    p_c2_d.font.size = Pt(12)
    p_c2_d.font.color.rgb = TEXT_MUTED

    # Card 3: Isolation of Info
    card3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.6), Inches(3.6), Inches(3.8))
    card3.fill.solid()
    card3.fill.fore_color.rgb = CARD_BG
    card3.line.color.rgb = CARD_BG
    tf_c3 = card3.text_frame
    tf_c3.word_wrap = True
    tf_c3.margin_left = tf_c3.margin_right = tf_c3.margin_top = tf_c3.margin_bottom = Inches(0.3)
    p_c3_t = tf_c3.paragraphs[0]
    p_c3_t.text = "🔍 Information Silos"
    p_c3_t.font.name = FONT_TITLE
    p_c3_t.font.size = Pt(18)
    p_c3_t.font.bold = True
    p_c3_t.font.color.rgb = ACCENT_ORANGE
    p_c3_t.space_after = Pt(14)
    p_c3_d = tf_c3.add_paragraph()
    p_c3_d.text = "Knowing *who* is sitting *where* or *which* floor a laboratory is on requires asking multiple staff. No centralized, real-time database exists for campus locations."
    p_c3_d.font.name = FONT_BODY
    p_c3_d.font.size = Pt(12)
    p_c3_d.font.color.rgb = TEXT_MUTED

    add_footer(slide2, 2)

    # ==================== SLIDE 3: OUR IDEA / VALUE PROP ====================
    slide3 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide3, BG_COLOR)
    add_header(slide3, "Our Idea & Mission")
    
    # Left Block: Large Quote style
    quote_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.2))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = CARD_BG
    quote_card.line.color.rgb = ACCENT_BLUE
    quote_card.line.width = Pt(1.5)
    tf_q = quote_card.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = Inches(0.4)
    
    p_q_open = tf_q.paragraphs[0]
    p_q_open.text = "THE VISION"
    p_q_open.font.name = FONT_BODY
    p_q_open.font.size = Pt(11)
    p_q_open.font.bold = True
    p_q_open.font.color.rgb = ACCENT_BLUE
    p_q_open.space_after = Pt(16)

    p_q_main = tf_q.add_paragraph()
    p_q_main.text = "“Our goal is to make finding any place inside the campus as simple as searching for a place on a map.”"
    p_q_main.font.name = FONT_TITLE
    p_q_main.font.size = Pt(24)
    p_q_main.font.bold = True
    p_q_main.font.italic = True
    p_q_main.font.color.rgb = TEXT_MAIN
    p_q_main.space_after = Pt(20)
    
    p_q_s = tf_q.add_paragraph()
    p_q_s.text = "- SITCOE Smart Campus Navigation System"
    p_q_s.font.name = FONT_BODY
    p_q_s.font.size = Pt(12)
    p_q_s.font.color.rgb = ACCENT_ORANGE

    # Right Content: Concept Details
    right_box = slide3.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.2))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "A Unified Navigation Hub"
    p_r1.font.name = FONT_TITLE
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = ACCENT_ORANGE
    p_r1.space_after = Pt(8)
    
    p_r1_d = tf_r.add_paragraph()
    p_r1_d.text = "We propose an intelligent, web-based platform tailored for SITCOE. It integrates search capabilities, directory details, and interactive geographical routing into a single seamless interface."
    p_r1_d.font.name = FONT_BODY
    p_r1_d.font.size = Pt(13)
    p_r1_d.font.color.rgb = TEXT_MUTED
    p_r1_d.space_after = Pt(24)
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "Solving the Last-Mile Indoor Challenge"
    p_r2.font.name = FONT_TITLE
    p_r2.font.size = Pt(20)
    p_r2.font.bold = True
    p_r2.font.color.rgb = ACCENT_ORANGE
    p_r2.space_after = Pt(8)
    
    p_r2_d = tf_r.add_paragraph()
    p_r2_d.text = "Standard tools like Google Maps stop at the campus boundary. Our solution map-routes the interior: linking halls, labs, administrative offices, and departments at a high resolution."
    p_r2_d.font.name = FONT_BODY
    p_r2_d.font.size = Pt(13)
    p_r2_d.font.color.rgb = TEXT_MUTED

    add_footer(slide3, 3)

    # ==================== SLIDE 4: HOW IT WORKS ====================
    slide4 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide4, BG_COLOR)
    add_header(slide4, "How It Works")
    
    steps = [
        ("Step 1: Input", "🎙️ Search / Voice", "User types a destination (e.g. 'CSE Room 38') or taps the microphone to speak their command naturally."),
        ("Step 2: Processing", "🤖 Location Match", "The system parses inputs in real-time, matching queries against a campus directory of departments, rooms, and faculty offices."),
        ("Step 3: Visualization", "🗺️ Interactive Map", "The digital campus map plots the route instantly, placing a custom marker and drawing the suggested paths overlay."),
        ("Step 4: Arrival", "🚶 Follow & Navigate", "The user receives walking guidance directly to the destination point, eliminating confusion and wrong turns.")
    ]
    
    for i, (step_num, title, desc) in enumerate(steps):
        left_pos = Inches(0.8 + i * 3.0)
        
        # Step header box
        s_box = slide4.shapes.add_textbox(left_pos, Inches(2.2), Inches(2.7), Inches(0.4))
        p_s = s_box.text_frame.paragraphs[0]
        p_s.text = step_num
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_BLUE
        
        # Step Card
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.7), Inches(2.7), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_ORANGE if i == 2 else CARD_BG # Highlight the map step
        card.line.width = Pt(1.5)
        
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = Inches(0.25)
        
        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(12)
        
        p_d = tf_c.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide4, 4)

    # ==================== SLIDE 5: MAIN FEATURES ====================
    slide5 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide5, BG_COLOR)
    add_header(slide5, "Main Features")
    
    # 2x3 Grid of features
    features = [
        ("🔎 Smart Search", "Instant fuzzy-search that filters campus spaces, faculties, and utility rooms as you type."),
        ("🎙️ Voice Navigation", "Hands-free speech commands parsed directly using the browser's Web Speech APIs."),
        ("🗺️ Interactive Map", "A dynamic GIS-powered campus canvas highlighting specific buildings and nodes."),
        ("🏫 Dept-Wise Workspace", "Filterable workspaces focusing on CSE, MECH, CIVIL, Electrical, ECE, AI&DS, and Admin."),
        ("📚 Room & Lab Library", "Granular listings details including room capacities, lab types, and equipment designations."),
        ("👨‍🏫 Faculty Directory", "Find staff workspace details, designations, departments, and map locations immediately.")
    ]
    
    for idx, (title, desc) in enumerate(features):
        row = idx // 3
        col = idx % 3
        left = Inches(0.8 + col * 3.9)
        top = Inches(2.2 + row * 2.3)
        
        feat_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.0))
        feat_card.fill.solid()
        feat_card.fill.fore_color.rgb = CARD_BG
        feat_card.line.color.rgb = CARD_BG
        
        tf_f = feat_card.text_frame
        tf_f.word_wrap = True
        tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = Inches(0.2)
        
        p_t = tf_f.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE
        p_t.space_after = Pt(6)
        
        p_d = tf_f.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide5, 5)

    # ==================== SLIDE 6: INNOVATION ====================
    slide6 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide6, BG_COLOR)
    add_header(slide6, "Key Innovation")
    
    # Left card: Navigation + Directory
    c_left = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = ACCENT_ORANGE
    c_left.line.width = Pt(1.5)
    tf_l = c_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = Inches(0.4)
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "🎯 Unified Hybrid Platform"
    p_l1.font.name = FONT_TITLE
    p_l1.font.size = Pt(20)
    p_l1.font.bold = True
    p_l1.font.color.rgb = TEXT_MAIN
    p_l1.space_after = Pt(12)
    
    p_l2 = tf_l.add_paragraph()
    p_l2.text = "We combine detailed navigation maps directly with directory databases.\n\nWhile standard solutions focus either solely on listing data (faculty directories) or solely on basic maps (Google Maps), our application binds them dynamically. Tapping a faculty member's name instantly plots their cabin coordinates on the map."
    p_l2.font.name = FONT_BODY
    p_l2.font.size = Pt(13)
    p_l2.font.color.rgb = TEXT_MUTED

    # Right cards (2 vertical blocks)
    c_r1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.0))
    c_r1.fill.solid()
    c_r1.fill.fore_color.rgb = CARD_BG
    c_r1.line.color.rgb = CARD_BG
    tf_r1 = c_r1.text_frame
    tf_r1.word_wrap = True
    tf_r1.margin_left = tf_r1.margin_right = tf_r1.margin_top = tf_r1.margin_bottom = Inches(0.3)
    p_r1_t = tf_r1.paragraphs[0]
    p_r1_t.text = "📍 College Campus Optimizations"
    p_r1_t.font.name = FONT_TITLE
    p_r1_t.font.size = Pt(16)
    p_r1_t.font.bold = True
    p_r1_t.font.color.rgb = ACCENT_BLUE
    p_r1_t.space_after = Pt(6)
    p_r1_d = tf_r1.add_paragraph()
    p_r1_d.text = "Features customized filters for college-specific constructs, like classroom blocks, laboratory schedules, and departmental partitions."
    p_r1_d.font.name = FONT_BODY
    p_r1_d.font.size = Pt(11.5)
    p_r1_d.font.color.rgb = TEXT_MUTED

    c_r2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.0))
    c_r2.fill.solid()
    c_r2.fill.fore_color.rgb = CARD_BG
    c_r2.line.color.rgb = CARD_BG
    tf_r2 = c_r2.text_frame
    tf_r2.word_wrap = True
    tf_r2.margin_left = tf_r2.margin_right = tf_r2.margin_top = tf_r2.margin_bottom = Inches(0.3)
    p_r2_t = tf_r2.paragraphs[0]
    p_r2_t.text = "📶 Foundation for Indoor Positioning"
    p_r2_t.font.name = FONT_TITLE
    p_r2_t.font.size = Pt(16)
    p_r2_t.font.bold = True
    p_r2_t.font.color.rgb = ACCENT_BLUE
    p_r2_t.space_after = Pt(6)
    p_r2_d = tf_r2.add_paragraph()
    p_r2_d.text = "Engineered with coordinate hooks ready to integrate with BLE beacons, Wi-Fi RTT nodes, or QR-code based room synchronization."
    p_r2_d.font.name = FONT_BODY
    p_r2_d.font.size = Pt(11.5)
    p_r2_d.font.color.rgb = TEXT_MUTED

    add_footer(slide6, 6)

    # ==================== SLIDE 7: EXPECTED BENEFITS ====================
    slide7 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide7, BG_COLOR)
    add_header(slide7, "Expected Benefits")
    
    benefits = [
        ("⚡ Extreme Time Saving", "Locates exam halls and faculty cabins in seconds, completely avoiding late-arrival stress during critical academic hours."),
        ("😌 Reduced Confusion", "Clear, visualized paths remove reliance on vague verbal directions like 'go straight then take the second left'."),
        ("🤝 Enhanced Visitor Experience", "Welcoming environment for parents, visiting delegates, recruiters, and workshop attendees navigating SITCOE."),
        ("♿ Accessibility Boost", "Can map and prioritize wheelchair-friendly ramps, escalators, and elevator lobbies for differently-abled users.")
    ]
    
    for idx, (title, desc) in enumerate(benefits):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(2.2 + row * 2.3)
        
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        
        tf_b = card.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.3)
        
        p_t = tf_b.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_ORANGE
        p_t.space_after = Pt(6)
        
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide7, 7)

    # ==================== SLIDE 8: FUTURE SCOPE ====================
    slide8 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide8, BG_COLOR)
    add_header(slide8, "Future Scope & Roadmap")
    
    scopes = [
        ("📍 Indoor Positioning", "Integrating BLE beacons or smartphone magnetic compass maps for real-time blue-dot positioning inside closed buildings."),
        ("🗣️ Turn-by-Turn Audio Guide", "Adding real-time audio guidance (e.g. 'Turn right at computer lab') for visually impaired users."),
        ("🚨 Emergency Route Planner", "Dynamic evacuation pathways that recalculate safe exit routes when hazard reports are broadcasted."),
        ("📱 Native Mobile App", "Launching Android & iOS applications leveraging native device APIs for seamless geolocation and background services.")
    ]
    
    for idx, (title, desc) in enumerate(scopes):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(2.2 + row * 2.3)
        
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        
        tf_s = card.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = Inches(0.3)
        
        p_t = tf_s.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE
        p_t.space_after = Pt(6)
        
        p_d = tf_s.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide8, 8)

    # ==================== SLIDE 9: CONCLUSION ====================
    slide9 = prs.slides.add_slide(slide_layout)
    apply_solid_background(slide9, BG_COLOR)
    add_header(slide9, "Conclusion & Next Steps", category="SUMMARY")
    
    # Large conclusion slide main highlight
    center_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.2))
    center_card.fill.solid()
    center_card.fill.fore_color.rgb = CARD_BG
    center_card.line.color.rgb = ACCENT_BLUE
    center_card.line.width = Pt(1.5)
    
    tf_conc = center_card.text_frame
    tf_conc.word_wrap = True
    tf_conc.margin_left = tf_conc.margin_right = tf_conc.margin_top = tf_conc.margin_bottom = Inches(0.5)
    
    p_c_main = tf_conc.paragraphs[0]
    p_c_main.text = "Making Campus Navigation Simple and Accessible"
    p_c_main.font.name = FONT_TITLE
    p_c_main.font.size = Pt(24)
    p_c_main.font.bold = True
    p_c_main.font.color.rgb = ACCENT_ORANGE
    p_c_main.space_after = Pt(16)
    
    p_c_body = tf_conc.add_paragraph()
    p_c_body.text = "The SITCOE Smart Campus Navigation System bridges the gap between complex physical geography and administrative information systems. By integrating a dynamic visual layout with smart search, voice commands, and departmental directory workspaces, we provide a modern asset for visitors, faculty, and students."
    p_c_body.font.name = FONT_BODY
    p_c_body.font.size = Pt(13)
    p_c_body.font.color.rgb = TEXT_MUTED
    p_c_body.space_after = Pt(24)
    
    p_q = tf_conc.add_paragraph()
    p_q.text = "Thank You!   |   Any Questions?"
    p_q.font.name = FONT_TITLE
    p_q.font.size = Pt(20)
    p_q.font.bold = True
    p_q.font.color.rgb = TEXT_MAIN
    
    add_footer(slide9, 9)

    # Save
    output_filename = "SITCOE_Smart_Campus_Navigation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_presentation()
